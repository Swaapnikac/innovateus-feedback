"""Export service — raw, structured, and user-testing CSVs + PDF/PPTX.

Design goals for the Apr 22 soft launch:

* Every important field is a dedicated column (no nested JSON blobs for
  fields we care about).
* Submission-level rollups come from ``metrics_service.compute_submission_rollup``
  so the CSVs and the dashboard always agree.
* Three CSVs only:
  - ``raw.csv`` — one row per submission, per-question block repeated for
    each survey question, with dedicated followup_1 / followup_2 columns.
  - ``structured.csv`` — one row per main question response, follow-ups in
    columns, submission-level rollups repeated on every row.
  - ``user-testing.csv`` — one row per submission, hypothesis flags included.

All three CSVs call a single "bundle" builder that reads ORM objects plus a
``SubmissionRollup`` and returns a flat dict. Consistency is enforced by
having every CSV read from that same bundle.
"""
from __future__ import annotations

import csv
import io
import json
import uuid
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable, Optional

from app.models import Answer, Cohort, Extraction, ExtractionReview, Submission
from app.services.metrics_service import (
    OPEN_TYPES,
    SubmissionRollup,
    compute_hypothesis_flags,
    compute_submission_rollup,
    word_count,
)


SURVEY_CONFIG_PATH = Path(__file__).resolve().parents[4] / "docs" / "survey-config" / "survey-en.json"

QUALTRICS_PAYLOAD_VERSION = "v1.2024-04"
CONSENT_VERSION = "1.0"
PRIVACY_NOTICE_VERSION = "1.0"


# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────


def _load_default_survey() -> dict:
    with open(SURVEY_CONFIG_PATH, "r", encoding="utf-8") as f:
        return json.load(f)


def _questions_from_config(survey_config: Optional[dict]) -> list[dict]:
    """Return a list of question dicts in survey order."""
    config = survey_config if survey_config else _load_default_survey()
    return list(config.get("questions", []) or [])


def _iso(ts) -> str:
    if ts is None:
        return ""
    if isinstance(ts, datetime):
        return ts.isoformat()
    return str(ts)


def _bool_to_str(v) -> str:
    if v is None:
        return ""
    return "true" if v else "false"


def _join_list(v) -> str:
    if not v:
        return ""
    if isinstance(v, list):
        return " | ".join(str(x) for x in v)
    return str(v)


def _format_answer(answer_raw: Optional[str], question_type: Optional[str]) -> str:
    if not answer_raw:
        return ""
    if question_type == "multi":
        try:
            items = json.loads(answer_raw)
            if isinstance(items, list):
                return " | ".join(str(item) for item in items)
        except (json.JSONDecodeError, TypeError):
            pass
    if question_type == "matrix":
        try:
            parsed = json.loads(answer_raw)
            if isinstance(parsed, dict):
                return " | ".join(f"{k}: {v}" for k, v in parsed.items())
        except (json.JSONDecodeError, TypeError):
            pass
    if question_type == "ranking":
        try:
            parsed = json.loads(answer_raw)
            if isinstance(parsed, list):
                return " > ".join(str(item) for item in parsed)
        except (json.JSONDecodeError, TypeError):
            pass
    if question_type == "yesno":
        return "Yes" if answer_raw == "yes" else "No" if answer_raw == "no" else answer_raw
    return answer_raw


# ─────────────────────────────────────────────────────────────────────────────
# Bundle: everything we need to know about one submission, flattened
# ─────────────────────────────────────────────────────────────────────────────


@dataclass
class SubmissionBundle:
    submission: Submission
    answers: list[Answer]
    extraction: Optional[Extraction]
    review: Optional[ExtractionReview]
    cohort: Optional[Cohort]
    rollup: SubmissionRollup
    hypothesis_flags: dict


def build_bundles(
    submissions: list[Submission],
    answers_by_sub: dict[str, list[Answer]],
    extractions_by_sub: dict[str, Extraction] | None = None,
    reviews_by_sub: dict[str, ExtractionReview] | None = None,
    cohorts_by_id: dict[str, Cohort] | None = None,
) -> list[SubmissionBundle]:
    extractions_by_sub = extractions_by_sub or {}
    reviews_by_sub = reviews_by_sub or {}
    cohorts_by_id = cohorts_by_id or {}
    out: list[SubmissionBundle] = []
    for s in submissions:
        sid = str(s.id)
        ans = answers_by_sub.get(sid, [])
        rollup = compute_submission_rollup(s, ans)
        review = reviews_by_sub.get(sid)
        flags = compute_hypothesis_flags(s, rollup, review)
        out.append(
            SubmissionBundle(
                submission=s,
                answers=list(ans),
                extraction=extractions_by_sub.get(sid),
                review=review,
                cohort=cohorts_by_id.get(str(s.cohort_id)),
                rollup=rollup,
                hypothesis_flags=flags,
            )
        )
    return out


# ─────────────────────────────────────────────────────────────────────────────
# Column builders — produce the named-column dicts each CSV needs
# ─────────────────────────────────────────────────────────────────────────────


def _session_columns(bundle: SubmissionBundle) -> dict:
    sub = bundle.submission
    cohort = bundle.cohort
    return {
        "submission_id": str(sub.id),
        "session_id": sub.ip_hash[:16] if sub.ip_hash else "",
        "participant_anonymous_id": sub.ip_hash[:8] if sub.ip_hash else "",
        "survey_id": str(sub.cohort_id),
        "survey_name": cohort.name if cohort else "",
        "survey_version": sub.survey_version or "",
        "survey_version_created_at": "",
        "program_id": str(sub.cohort_id),
        "program_name": cohort.course_name if cohort else "",
        "program_type": (cohort.program_type if cohort else "") or "",
        "cohort_id": str(sub.cohort_id),
        "facilitator_id": cohort.facilitator_email if cohort else "",
        "facilitator_name": cohort.facilitator_name if cohort else "",
        "launch_phase": (cohort.launch_phase if cohort else "") or "soft_launch",
        "source_channel": (cohort.source_channel if cohort else "") or "",
        "survey_link_id": "",
        "qr_code_id": "",
        "started_at": _iso(sub.created_at),
        "completed_at": _iso(sub.completed_at),
        "abandoned_at": _iso(sub.abandoned_at),
        "last_activity_at": _iso(sub.last_activity_at),
        "submission_status": sub.status or "",
        "total_time_to_complete_sec": sub.time_to_complete_sec if sub.time_to_complete_sec is not None else "",
        "completed_flag": _bool_to_str(sub.status == "completed"),
        "abandoned_flag": _bool_to_str(sub.status == "abandoned"),
        "partial_flag": _bool_to_str(sub.status not in ("completed", "abandoned") and bundle.rollup.total_questions_answered > 0),
        "final_completion_percent": round(
            bundle.rollup.total_questions_answered / max(1, bundle.rollup.total_questions_seen) * 100, 2
        ) if bundle.rollup.total_questions_seen else 0,
    }


def _device_columns(bundle: SubmissionBundle) -> dict:
    sub = bundle.submission
    return {
        "browser_name": sub.browser_name or "",
        "browser_version": sub.browser_version or "",
        "os_name": sub.os_name or "",
        "os_version": sub.os_version or "",
        "device_type": sub.device_type or "",
        "user_agent": sub.user_agent or "",
        "screen_size": sub.screen_size or "",
        "microphone_permission_status": sub.mic_permission_status or "",
        "microphone_permission_prompted_at": _iso(sub.mic_permission_prompted_at),
        "voice_supported_in_browser": _bool_to_str(sub.voice_supported_in_browser),
        "speech_recognition_provider_used": sub.speech_recognition_provider_used or "",
        "fallback_provider_used": "",
        "network_status_summary": sub.connection_type or "",
        "connection_type": sub.connection_type or "",
        "page_load_time_ms": sub.page_load_time_ms if sub.page_load_time_ms is not None else "",
        "avg_api_latency_ms": sub.avg_api_latency_ms if sub.avg_api_latency_ms is not None else "",
        "max_api_latency_ms": sub.max_api_latency_ms if sub.max_api_latency_ms is not None else "",
        "total_api_failures": sub.total_api_failures or 0,
        "total_retry_count": 0,
        "sentry_error_count": sub.sentry_error_count or 0,
        "client_error_count": sub.client_error_count or 0,
        "critical_error_flag": _bool_to_str(sub.critical_error_flag),
        "timeout_count": sub.timeout_count or 0,
    }


def _question_columns(answer: Optional[Answer], question: dict, order: int) -> dict:
    q_id = question.get("id") or ""
    q_text = question.get("text") or ""
    q_type = question.get("type") or ""
    q_category = "OPEN" if q_type.lower() in {"open", "voice"} else "CLOSE"
    cols: dict = {
        "question_id": q_id,
        "question_text": q_text,
        "question_type": q_type,
        # OPEN / CLOSE — capitalised analyst-facing marker. Sits next to
        # ``question_type`` (which carries the technical type like
        # ``rating`` / ``mcq`` / ``open``) so both audiences are served
        # without renaming the existing column.
        "question_category": q_category,
        "question_order": order,
        "section_name": question.get("group") or "",
        "block_name": question.get("group") or "",
        "is_required": _bool_to_str(bool(question.get("required"))),
    }
    if answer is None:
        cols.update(
            {
                "response_received": _bool_to_str(False),
                "answer_raw": "",
                "answer_cleaned": "",
                "answer_word_count": 0,
                "answer_char_count": 0,
                "answer_language": "",
                "answer_input_mode": "",
                "answer_started_at": "",
                "answer_completed_at": "",
                "answer_duration_sec": "",
                "answer_skipped": _bool_to_str(True),
                "skip_reason": "",
                "changed_answer_flag": _bool_to_str(False),
                "edited_after_transcription_flag": _bool_to_str(False),
                "voice_used_flag": _bool_to_str(False),
                "voice_started_at": "",
                "voice_ended_at": "",
                "voice_duration_sec": "",
                "speech_to_text_provider": "",
                "transcript_raw": "",
                "transcript_final": "",
                "transcript_edit_distance": "",
                "user_edited_transcript_flag": _bool_to_str(False),
                "silence_auto_stop_triggered_flag": _bool_to_str(False),
                "audio_capture_error_flag": _bool_to_str(False),
                "microphone_denied_flag": _bool_to_str(False),
                "microphone_denied_stage": "",
                "switched_from_voice_to_text_flag": _bool_to_str(False),
                "switch_from_voice_to_text_stage": "",
                "switch_reason": "",
                "switched_from_text_to_voice_flag": _bool_to_str(False),
                "initial_response_is_vague_flag": "",
                "vagueness_score_initial": "",
                "vagueness_reason_initial": "",
                "followup_1_shown_flag": _bool_to_str(False),
                "followup_1_question": "",
                "followup_1_answer": "",
                "followup_1_answer_word_count": 0,
                "followup_1_input_mode": "",
                "followup_1_is_vague_flag": "",
                "followup_1_vagueness_score": "",
                "followup_1_improved_specificity_flag": "",
                "followup_1_skipped_flag": _bool_to_str(False),
                "followup_2_shown_flag": _bool_to_str(False),
                "followup_2_question": "",
                "followup_2_answer": "",
                "followup_2_answer_word_count": 0,
                "followup_2_input_mode": "",
                "followup_2_is_vague_flag": "",
                "followup_2_vagueness_score": "",
                "followup_2_improved_specificity_flag": "",
                "followup_2_skipped_flag": _bool_to_str(False),
                "final_response_specific_flag": "",
                "total_response_word_count_after_followups": 0,
                "specificity_improved_after_followups_flag": "",
                "total_followups_shown_for_question": 0,
                "total_followups_answered_for_question": 0,
            }
        )
        return cols

    raw = answer.answer_raw or ""
    cleaned = _format_answer(raw, answer.question_type)
    voice_flag = (answer.input_mode or "") == "voice"
    followups_shown = (1 if answer.followup_1 else 0) + (1 if answer.followup_2 else 0)
    followups_answered = (1 if answer.followup_1_answer else 0) + (1 if answer.followup_2_answer else 0)
    fu1_wc = answer.followup_1_word_count if answer.followup_1_word_count is not None else word_count(answer.followup_1_answer or "")
    fu2_wc = answer.followup_2_word_count if answer.followup_2_word_count is not None else word_count(answer.followup_2_answer or "")
    answer_wc = answer.answer_word_count if answer.answer_word_count is not None else word_count(raw)
    total_wc_after = answer_wc + fu1_wc + fu2_wc

    cols.update(
        {
            "response_received": _bool_to_str(bool(raw.strip())),
            "answer_raw": raw,
            "answer_cleaned": cleaned,
            "answer_word_count": answer_wc,
            "answer_char_count": answer.answer_char_count if answer.answer_char_count is not None else len(raw),
            "answer_language": answer.answer_language or "",
            "answer_input_mode": answer.input_mode or "",
            "answer_started_at": _iso(answer.answer_started_at),
            "answer_completed_at": _iso(answer.answer_completed_at),
            "answer_duration_sec": answer.answer_duration_sec if answer.answer_duration_sec is not None else "",
            "answer_skipped": _bool_to_str(bool(answer.answer_skipped)),
            "skip_reason": answer.skip_reason or "",
            "changed_answer_flag": _bool_to_str(bool(answer.changed_answer_flag)),
            "edited_after_transcription_flag": _bool_to_str(bool(answer.edited_after_transcription_flag)),
            "voice_used_flag": _bool_to_str(voice_flag),
            "voice_started_at": _iso(answer.answer_started_at) if voice_flag else "",
            "voice_ended_at": _iso(answer.answer_completed_at) if voice_flag else "",
            "voice_duration_sec": answer.voice_duration_sec if answer.voice_duration_sec is not None else "",
            "speech_to_text_provider": "whisper" if voice_flag else "",
            "transcript_raw": answer.transcript_raw or "",
            "transcript_final": answer.transcript_final or (cleaned if voice_flag else ""),
            "transcript_edit_distance": answer.transcript_edit_distance if answer.transcript_edit_distance is not None else "",
            "user_edited_transcript_flag": _bool_to_str(bool(answer.user_edited_transcript_flag)),
            "silence_auto_stop_triggered_flag": _bool_to_str(bool(answer.silence_auto_stop_triggered_flag)),
            "audio_capture_error_flag": _bool_to_str(bool(answer.audio_capture_error_flag)),
            "microphone_denied_flag": _bool_to_str(bool(answer.mic_denied_flag)),
            "microphone_denied_stage": answer.question_id if answer.mic_denied_flag else "",
            "switched_from_voice_to_text_flag": _bool_to_str(bool(answer.switched_from_voice_to_text_flag)),
            "switch_from_voice_to_text_stage": answer.question_id if answer.switched_from_voice_to_text_flag else "",
            "switch_reason": "",
            "switched_from_text_to_voice_flag": _bool_to_str(bool(answer.switched_from_text_to_voice_flag)),
            "initial_response_is_vague_flag": _bool_to_str(answer.is_vague) if answer.is_vague is not None else "",
            "vagueness_score_initial": answer.vagueness_score_initial if answer.vagueness_score_initial is not None else "",
            "vagueness_reason_initial": answer.vagueness_reason_initial or "",
            "followup_1_shown_flag": _bool_to_str(bool(answer.followup_1)),
            "followup_1_question": answer.followup_1 or "",
            "followup_1_answer": answer.followup_1_answer or "",
            "followup_1_answer_word_count": fu1_wc,
            "followup_1_input_mode": answer.followup_1_input_mode or "",
            "followup_1_is_vague_flag": _bool_to_str(answer.followup_1_is_vague) if answer.followup_1_is_vague is not None else "",
            "followup_1_vagueness_score": answer.followup_1_vagueness_score if answer.followup_1_vagueness_score is not None else "",
            "followup_1_improved_specificity_flag": _bool_to_str(
                answer.is_vague and answer.followup_1_is_vague is False
            ) if (answer.followup_1_answer and answer.is_vague is not None) else "",
            "followup_1_skipped_flag": _bool_to_str(bool(answer.followup_1_skipped_flag)),
            "followup_2_shown_flag": _bool_to_str(bool(answer.followup_2)),
            "followup_2_question": answer.followup_2 or "",
            "followup_2_answer": answer.followup_2_answer or "",
            "followup_2_answer_word_count": fu2_wc,
            "followup_2_input_mode": answer.followup_2_input_mode or "",
            "followup_2_is_vague_flag": _bool_to_str(answer.followup_2_is_vague) if answer.followup_2_is_vague is not None else "",
            "followup_2_vagueness_score": answer.followup_2_vagueness_score if answer.followup_2_vagueness_score is not None else "",
            "followup_2_improved_specificity_flag": _bool_to_str(
                answer.followup_1_is_vague and answer.followup_2_is_vague is False
            ) if (answer.followup_2_answer and answer.followup_1_is_vague is not None) else "",
            "followup_2_skipped_flag": _bool_to_str(bool(answer.followup_2_skipped_flag)),
            "final_response_specific_flag": _bool_to_str(answer.final_response_specific_flag) if answer.final_response_specific_flag is not None else "",
            "total_response_word_count_after_followups": total_wc_after,
            "specificity_improved_after_followups_flag": _bool_to_str(answer.specificity_improved_after_followups_flag) if answer.specificity_improved_after_followups_flag is not None else "",
            "total_followups_shown_for_question": followups_shown,
            "total_followups_answered_for_question": followups_answered,
        }
    )
    return cols


def _extraction_columns(bundle: SubmissionBundle) -> dict:
    ext = bundle.extraction
    if ext is None:
        return {
            "extraction_model_name": "",
            "extraction_model_version": "",
            "extraction_prompt_version": "",
            "extraction_run_at": "",
            "extraction_success_flag": "",
            "extracted_summary": "",
            "extracted_themes": "",
            "extracted_barriers": "",
            "extracted_enablers": "",
            "extracted_success_story": "",
            "extracted_sentiment": "",
            "extracted_action_items": "",
            "extraction_confidence": "",
            "extraction_error_flag": "",
            "extraction_error_message": "",
        }
    summary_parts = [
        p for p in [ext.what_was_tried, ext.outcome_or_expected_outcome, ext.public_benefit] if p
    ]
    return {
        "extraction_model_name": ext.model_name or "",
        "extraction_model_version": ext.model_version or "",
        "extraction_prompt_version": ext.prompt_version or "",
        "extraction_run_at": _iso(ext.run_at),
        "extraction_success_flag": _bool_to_str(ext.success_flag) if ext.success_flag is not None else "",
        "extracted_summary": " | ".join(summary_parts),
        "extracted_themes": _join_list(ext.top_themes),
        "extracted_barriers": _join_list(ext.barriers),
        "extracted_enablers": _join_list(ext.enablers),
        "extracted_success_story": ext.success_story_candidate or "",
        "extracted_sentiment": ext.sentiment or "",
        "extracted_action_items": ext.planned_task_or_workflow or "",
        "extraction_confidence": ext.confidence if ext.confidence is not None else "",
        "extraction_error_flag": _bool_to_str(bool(ext.error_message)),
        "extraction_error_message": ext.error_message or "",
    }


def _review_columns(bundle: SubmissionBundle) -> dict:
    rev = bundle.review
    completed = bundle.submission.status == "completed"
    return {
        "review_required_flag": _bool_to_str(completed),
        "reviewed_submission_flag": _bool_to_str(bool(rev)),
        "reviewed_by": rev.reviewed_by if rev else "",
        "reviewed_at": _iso(rev.reviewed_at) if rev else "",
        "extraction_accuracy_rating": rev.accuracy_rating if (rev and rev.accuracy_rating is not None) else "",
        "extraction_usefulness_rating": rev.usefulness_rating if (rev and rev.usefulness_rating is not None) else "",
        "extraction_useful_flag": _bool_to_str(rev.useful_flag) if (rev and rev.useful_flag is not None) else "",
        "extraction_accuracy_notes": rev.accuracy_notes if rev else "",
        "extraction_usefulness_notes": rev.usefulness_notes if rev else "",
    }


def _qualtrics_columns(bundle: SubmissionBundle) -> dict:
    sub = bundle.submission
    attempted = (sub.qualtrics_sync_attempt_count or 0) > 0
    succeeded = sub.qualtrics_synced_at is not None
    status = "success" if succeeded else ("failed" if attempted else "pending")
    latency_sec = (sub.qualtrics_sync_latency_ms or 0) / 1000.0 if sub.qualtrics_sync_latency_ms else ""
    return {
        "qualtrics_sync_attempted_flag": _bool_to_str(attempted),
        "qualtrics_sync_status": status,
        "qualtrics_sync_attempt_count": sub.qualtrics_sync_attempt_count or 0,
        "qualtrics_sync_first_attempt_at": _iso(sub.qualtrics_sync_first_attempt_at),
        "qualtrics_sync_last_attempt_at": _iso(sub.qualtrics_sync_last_attempt_at),
        "qualtrics_synced_at": _iso(sub.qualtrics_synced_at),
        "qualtrics_response_id": sub.qualtrics_response_id or "",
        "qualtrics_error_message": sub.qualtrics_sync_last_error or "",
        "qualtrics_payload_version": QUALTRICS_PAYLOAD_VERSION,
        "qualtrics_sync_success_flag": _bool_to_str(succeeded),
        "qualtrics_sync_latency_sec": latency_sec,
    }


def _feedback_columns(bundle: SubmissionBundle) -> dict:
    sub = bundle.submission
    return {
        "user_experience_rating": sub.experience_rating if sub.experience_rating is not None else "",
        "user_experience_text": sub.experience_feedback or "",
        "voice_experience_rating": sub.voice_experience_rating if sub.voice_experience_rating is not None else "",
        "voice_experience_text": sub.voice_experience_text or "",
        "confusion_flag": _bool_to_str(sub.confusion_flag) if sub.confusion_flag is not None else "",
        "confusion_step": sub.confusion_step or "",
        "would_use_again_flag": _bool_to_str(sub.would_use_again_flag) if sub.would_use_again_flag is not None else "",
        "preferred_mode_next_time": sub.preferred_mode_next_time or "",
        "reported_issue_flag": _bool_to_str(sub.reported_issue_flag) if sub.reported_issue_flag is not None else "",
        "reported_issue_text": sub.reported_issue_text or "",
    }


def _facilitator_columns(bundle: SubmissionBundle) -> dict:
    cohort = bundle.cohort
    if cohort is None:
        return {
            "facilitator_feedback_received_flag": _bool_to_str(False),
            "facilitator_feedback_text": "",
            "facilitator_reported_issue_flag": _bool_to_str(False),
            "facilitator_issue_type": "",
            "facilitator_issue_notes": "",
        }
    return {
        "facilitator_feedback_received_flag": _bool_to_str(bool(cohort.facilitator_feedback_text)),
        "facilitator_feedback_text": cohort.facilitator_feedback_text or "",
        "facilitator_reported_issue_flag": _bool_to_str(bool(cohort.facilitator_reported_issue_flag)),
        "facilitator_issue_type": cohort.facilitator_issue_type or "",
        "facilitator_issue_notes": cohort.facilitator_issue_notes or "",
    }


def _governance_columns(bundle: SubmissionBundle) -> dict:
    sub = bundle.submission
    return {
        "pii_detected_flag": _bool_to_str(sub.pii_detected_flag),
        "pii_redaction_applied_flag": _bool_to_str(sub.pii_redaction_applied_flag if sub.pii_redaction_applied_flag is not None else True),
        "consent_version": sub.consent_version or CONSENT_VERSION,
        "privacy_notice_version": sub.privacy_notice_version or PRIVACY_NOTICE_VERSION,
        "records_retention_tag": sub.records_retention_tag or "standard",
        "export_generated_at": datetime.now(timezone.utc).isoformat(),
    }


def _submission_rollup_columns(bundle: SubmissionBundle) -> dict:
    r = bundle.rollup
    return {
        "total_questions_seen": r.total_questions_seen,
        "total_questions_answered": r.total_questions_answered,
        "total_open_ended_questions_seen": r.total_open_ended_questions_seen,
        "total_open_ended_questions_answered": r.total_open_ended_questions_answered,
        "total_followups_shown": r.total_followups_shown,
        "total_followups_answered": r.total_followups_answered,
        "total_followups_skipped": r.total_followups_skipped,
        "voice_used_any_flag": _bool_to_str(r.voice_used_any_flag),
        "started_in_voice_flag": _bool_to_str(r.started_in_voice_flag) if r.started_in_voice_flag is not None else "",
        "ended_in_voice_flag": _bool_to_str(r.ended_in_voice_flag) if r.ended_in_voice_flag is not None else "",
        "switched_voice_to_text_any_flag": _bool_to_str(r.switched_voice_to_text_any_flag),
        "switched_text_to_voice_any_flag": _bool_to_str(r.switched_text_to_voice_any_flag),
        "voice_conversation_completed_flag": _bool_to_str(r.voice_conversation_completed_flag),
        "total_open_ended_word_count": r.total_open_ended_word_count,
        "avg_open_ended_word_count": r.avg_open_ended_word_count if r.avg_open_ended_word_count is not None else "",
        "avg_voice_open_ended_word_count": r.avg_voice_open_ended_word_count if r.avg_voice_open_ended_word_count is not None else "",
        "avg_text_open_ended_word_count": r.avg_text_open_ended_word_count if r.avg_text_open_ended_word_count is not None else "",
        "initial_vague_answer_count": r.initial_vague_answer_count,
        "vague_answer_count_after_followups": r.vague_answer_count_after_followups,
        "specificity_improvement_rate": r.specificity_improvement_rate if r.specificity_improvement_rate is not None else "",
    }


# ─────────────────────────────────────────────────────────────────────────────
# CSV generators
# ─────────────────────────────────────────────────────────────────────────────


def generate_raw_csv(
    bundles: list[SubmissionBundle],
    survey_config: Optional[dict] = None,
) -> str:
    """One row per submission, every field its own column.

    Per-question block (``q{i}_id``, ``q{i}_answer_raw``, etc.) is repeated in
    survey order for every question. Follow-ups have dedicated columns on the
    same row.
    """
    questions = _questions_from_config(survey_config)

    output = io.StringIO()
    writer = csv.writer(output)

    # Build static column list for a single submission
    session_cols = list(_session_columns(bundles[0]).keys()) if bundles else list(
        _session_columns(_empty_bundle_for_headers()).keys()
    )
    device_cols = list(_device_columns(bundles[0]).keys()) if bundles else list(
        _device_columns(_empty_bundle_for_headers()).keys()
    )

    # Per-question columns (prefixed by q1_OPEN_, q2_CLOSE_, ...). The
    # OPEN / CLOSE marker in the header makes it visible at a glance which
    # block is free-text vs constrained — handy when the CSV runs to
    # hundreds of columns and the survey question text is long.
    question_col_keys = list(
        _question_columns(None, {"id": "", "text": "", "type": "", "required": False}, 1).keys()
    )
    per_q_headers: list[str] = []
    for idx, q in enumerate(questions, 1):
        marker = "OPEN" if (q.get("type") or "").lower() in {"open", "voice"} else "CLOSE"
        prefix = f"q{idx}_{marker}_"
        per_q_headers.extend(prefix + k for k in question_col_keys)

    rollup_keys = list(_submission_rollup_columns(_empty_bundle_for_headers()).keys())
    ext_keys = list(_extraction_columns(_empty_bundle_for_headers()).keys())
    rev_keys = list(_review_columns(_empty_bundle_for_headers()).keys())
    qual_keys = list(_qualtrics_columns(_empty_bundle_for_headers()).keys())
    fb_keys = list(_feedback_columns(_empty_bundle_for_headers()).keys())
    fac_keys = list(_facilitator_columns(_empty_bundle_for_headers()).keys())
    gov_keys = list(_governance_columns(_empty_bundle_for_headers()).keys())

    headers = (
        session_cols
        + device_cols
        + per_q_headers
        + rollup_keys
        + ext_keys
        + rev_keys
        + qual_keys
        + fb_keys
        + fac_keys
        + gov_keys
    )
    writer.writerow(headers)

    for bundle in bundles:
        ans_by_qid = {a.question_id: a for a in bundle.answers}
        row: list = []
        for v in _session_columns(bundle).values():
            row.append(v)
        for v in _device_columns(bundle).values():
            row.append(v)
        for idx, q in enumerate(questions, 1):
            a = ans_by_qid.get(q.get("id", ""))
            cols = _question_columns(a, q, idx)
            for v in cols.values():
                row.append(v)
        for v in _submission_rollup_columns(bundle).values():
            row.append(v)
        for v in _extraction_columns(bundle).values():
            row.append(v)
        for v in _review_columns(bundle).values():
            row.append(v)
        for v in _qualtrics_columns(bundle).values():
            row.append(v)
        for v in _feedback_columns(bundle).values():
            row.append(v)
        for v in _facilitator_columns(bundle).values():
            row.append(v)
        for v in _governance_columns(bundle).values():
            row.append(v)
        writer.writerow(row)

    return output.getvalue()


def generate_structured_csv(
    bundles: list[SubmissionBundle],
    survey_config: Optional[dict] = None,
) -> str:
    """One row per main question response.

    Follow-ups live in dedicated columns on the same row. Submission-level
    rollups repeat on every row. This is the ONE structured CSV — do not add
    more structured CSVs.
    """
    questions = _questions_from_config(survey_config)
    q_order_by_id = {q.get("id"): i + 1 for i, q in enumerate(questions)}

    output = io.StringIO()
    writer = csv.writer(output)

    header_template = None
    rows: list[list] = []

    for bundle in bundles:
        session_cols = _session_columns(bundle)
        device_cols = _device_columns(bundle)
        rollup_cols = _submission_rollup_columns(bundle)
        ext_cols = _extraction_columns(bundle)
        rev_cols = _review_columns(bundle)
        qual_cols = _qualtrics_columns(bundle)
        fb_cols = _feedback_columns(bundle)
        fac_cols = _facilitator_columns(bundle)
        gov_cols = _governance_columns(bundle)

        # For each main question in the survey, emit one row (even if unanswered)
        ans_by_qid = {a.question_id: a for a in bundle.answers}
        for idx, q in enumerate(questions, 1):
            a = ans_by_qid.get(q.get("id", ""))
            q_cols = _question_columns(a, q, idx)

            row_dict: dict = {}
            row_dict.update(session_cols)
            row_dict.update(device_cols)
            row_dict.update(q_cols)
            row_dict.update(ext_cols)
            row_dict.update(rev_cols)
            row_dict.update(rollup_cols)
            row_dict.update(qual_cols)
            row_dict.update(fb_cols)
            row_dict.update(fac_cols)
            row_dict.update(gov_cols)

            if header_template is None:
                header_template = list(row_dict.keys())
                writer.writerow(header_template)
            rows.append([row_dict.get(k, "") for k in header_template])

    # Handle empty-submissions case — still write a header using an empty bundle
    if header_template is None:
        empty = _empty_bundle_for_headers()
        row_dict = {}
        row_dict.update(_session_columns(empty))
        row_dict.update(_device_columns(empty))
        row_dict.update(
            _question_columns(None, {"id": "", "text": "", "type": "", "required": False}, 1)
        )
        row_dict.update(_extraction_columns(empty))
        row_dict.update(_review_columns(empty))
        row_dict.update(_submission_rollup_columns(empty))
        row_dict.update(_qualtrics_columns(empty))
        row_dict.update(_feedback_columns(empty))
        row_dict.update(_facilitator_columns(empty))
        row_dict.update(_governance_columns(empty))
        writer.writerow(list(row_dict.keys()))

    for r in rows:
        writer.writerow(r)

    return output.getvalue()


def generate_user_testing_csv(bundles: list[SubmissionBundle]) -> str:
    """One row per submission, hypothesis-scored for the soft launch."""
    columns = [
        "submission_id",
        "session_id",
        "survey_id",
        "survey_name",
        "survey_version",
        "program_id",
        "program_name",
        "program_type",
        "participant_anonymous_id",
        "launch_phase",
        "started_at",
        "completed_at",
        "submission_status",
        "completed_flag",
        "abandoned_flag",
        "abandonment_stage",
        "completion_time_sec",
        "completion_time_min",
        "voice_used_any_flag",
        "voice_used_count",
        "text_used_count",
        "started_in_voice_flag",
        "ended_in_voice_flag",
        "switched_voice_to_text_flag",
        "switched_text_to_voice_flag",
        "voice_conversation_completed_flag",
        "total_open_ended_word_count",
        "avg_open_ended_word_count",
        "avg_voice_open_ended_word_count",
        "avg_text_open_ended_word_count",
        "initial_vague_answer_count",
        "vague_answer_count_after_followups",
        "specificity_improvement_rate",
        "followup_shown_count",
        "followup_answered_count",
        "followup_skip_count",
        "extraction_success_flag",
        "extraction_reviewed_flag",
        "extraction_useful_flag",
        "extraction_accuracy_rating",
        "extraction_usefulness_rating",
        "qualtrics_sync_success_flag",
        "qualtrics_sync_latency_sec",
        "critical_error_flag",
        "sentry_error_count",
        "client_error_count",
        "mic_permission_failure_flag",
        "browser_compatibility_issue_flag",
        "browser_name",
        "device_type",
        "participant_feedback_score",
        "participant_feedback_text",
        "voice_experience_score",
        "voice_experience_text",
        "confusion_flag",
        "would_use_again_flag",
        "h1_voice_more_detailed_support_flag",
        "h2_followups_improve_quality_support_flag",
        "h3_completion_success_flag",
        "h4_extraction_useful_support_flag",
        "h5_voice_natural_support_flag",
        "h6_device_compatibility_support_flag",
    ]

    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow(columns)

    for bundle in bundles:
        sub = bundle.submission
        r = bundle.rollup
        rev = bundle.review
        cohort = bundle.cohort
        flags = bundle.hypothesis_flags

        voice_count = sum(
            1 for a in bundle.answers if (a.input_mode or "") == "voice"
        )
        text_count = sum(1 for a in bundle.answers if (a.input_mode or "") == "text")
        mic_failure = sub.mic_permission_status == "denied"
        browser_issue = bool(sub.critical_error_flag) or mic_failure or (sub.total_api_failures or 0) > 0
        extraction_success = bool(bundle.extraction and bundle.extraction.success_flag)
        qualtrics_ok = bool(sub.qualtrics_synced_at)
        qualtrics_latency_sec = (sub.qualtrics_sync_latency_ms or 0) / 1000.0 if sub.qualtrics_sync_latency_ms else ""

        row = [
            str(sub.id),
            sub.ip_hash[:16] if sub.ip_hash else "",
            str(sub.cohort_id),
            cohort.name if cohort else "",
            sub.survey_version or "",
            str(sub.cohort_id),
            cohort.course_name if cohort else "",
            (cohort.program_type if cohort else "") or "",
            sub.ip_hash[:8] if sub.ip_hash else "",
            (cohort.launch_phase if cohort else "") or "soft_launch",
            _iso(sub.created_at),
            _iso(sub.completed_at),
            sub.status or "",
            _bool_to_str(sub.status == "completed"),
            _bool_to_str(sub.status == "abandoned"),
            sub.abandonment_stage or "",
            sub.time_to_complete_sec if sub.time_to_complete_sec is not None else "",
            round((sub.time_to_complete_sec or 0) / 60.0, 2) if sub.time_to_complete_sec else "",
            _bool_to_str(r.voice_used_any_flag),
            voice_count,
            text_count,
            _bool_to_str(r.started_in_voice_flag) if r.started_in_voice_flag is not None else "",
            _bool_to_str(r.ended_in_voice_flag) if r.ended_in_voice_flag is not None else "",
            _bool_to_str(r.switched_voice_to_text_any_flag),
            _bool_to_str(r.switched_text_to_voice_any_flag),
            _bool_to_str(r.voice_conversation_completed_flag),
            r.total_open_ended_word_count,
            r.avg_open_ended_word_count if r.avg_open_ended_word_count is not None else "",
            r.avg_voice_open_ended_word_count if r.avg_voice_open_ended_word_count is not None else "",
            r.avg_text_open_ended_word_count if r.avg_text_open_ended_word_count is not None else "",
            r.initial_vague_answer_count,
            r.vague_answer_count_after_followups,
            r.specificity_improvement_rate if r.specificity_improvement_rate is not None else "",
            r.total_followups_shown,
            r.total_followups_answered,
            r.total_followups_skipped,
            _bool_to_str(extraction_success),
            _bool_to_str(bool(rev)),
            _bool_to_str(rev.useful_flag) if (rev and rev.useful_flag is not None) else "",
            rev.accuracy_rating if (rev and rev.accuracy_rating is not None) else "",
            rev.usefulness_rating if (rev and rev.usefulness_rating is not None) else "",
            _bool_to_str(qualtrics_ok),
            qualtrics_latency_sec,
            _bool_to_str(bool(sub.critical_error_flag)),
            sub.sentry_error_count or 0,
            sub.client_error_count or 0,
            _bool_to_str(mic_failure),
            _bool_to_str(browser_issue),
            sub.browser_name or "",
            sub.device_type or "",
            sub.experience_rating if sub.experience_rating is not None else "",
            sub.experience_feedback or "",
            sub.voice_experience_rating if sub.voice_experience_rating is not None else "",
            sub.voice_experience_text or "",
            _bool_to_str(sub.confusion_flag) if sub.confusion_flag is not None else "",
            _bool_to_str(sub.would_use_again_flag) if sub.would_use_again_flag is not None else "",
            _bool_to_str(flags["h1_voice_more_detailed_support_flag"]) if flags["h1_voice_more_detailed_support_flag"] is not None else "",
            _bool_to_str(flags["h2_followups_improve_quality_support_flag"]) if flags["h2_followups_improve_quality_support_flag"] is not None else "",
            _bool_to_str(flags["h3_completion_success_flag"]) if flags["h3_completion_success_flag"] is not None else "",
            _bool_to_str(flags["h4_extraction_useful_support_flag"]) if flags["h4_extraction_useful_support_flag"] is not None else "",
            _bool_to_str(flags["h5_voice_natural_support_flag"]) if flags["h5_voice_natural_support_flag"] is not None else "",
            _bool_to_str(flags["h6_device_compatibility_support_flag"]) if flags["h6_device_compatibility_support_flag"] is not None else "",
        ]
        writer.writerow(row)

    return output.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# Qualtrics CSV — Qualtrics-importable shape, 3-row header.
#
# Same row set as raw/structured (one row per submission), but column shape
# matches what Qualtrics emits for this survey so the file round-trips back
# via Response Import. Open-ended answers + their AI follow-ups are merged
# into the main question column; follow-up text never gets its own column.
# ─────────────────────────────────────────────────────────────────────────────


def generate_qualtrics_csv(
    bundles: list[SubmissionBundle],
    target,  # qualtrics_mapper.ResolvedTarget
    survey_config: Optional[dict] = None,
) -> str:
    from app.services import qualtrics_mapper

    questions = _questions_from_config(survey_config)
    output = io.StringIO()
    writer = csv.writer(output)

    for header_row in qualtrics_mapper.build_qualtrics_csv_headers(questions, target):
        writer.writerow(header_row)

    for bundle in bundles:
        row = qualtrics_mapper.build_qualtrics_csv_row(
            submission=bundle.submission,
            answers=bundle.answers,
            cohort=bundle.cohort,
            questions=questions,
            target=target,
        )
        writer.writerow(row)

    return output.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# Placeholder bundle for building headers when there are no submissions
# ─────────────────────────────────────────────────────────────────────────────


def _empty_bundle_for_headers() -> SubmissionBundle:
    from uuid import uuid4

    sub = Submission(
        id=uuid4(),
        cohort_id=uuid4(),
        created_at=datetime.now(timezone.utc),
        status="started",
    )
    rollup = SubmissionRollup(
        submission_id=str(sub.id),
        total_questions_seen=0,
        total_questions_answered=0,
        total_open_ended_questions_seen=0,
        total_open_ended_questions_answered=0,
        total_followups_shown=0,
        total_followups_answered=0,
        total_followups_skipped=0,
        voice_used_any_flag=False,
        started_in_voice_flag=None,
        ended_in_voice_flag=None,
        switched_voice_to_text_any_flag=False,
        switched_text_to_voice_any_flag=False,
        voice_conversation_completed_flag=False,
        total_open_ended_word_count=0,
        avg_open_ended_word_count=None,
        avg_voice_open_ended_word_count=None,
        avg_text_open_ended_word_count=None,
        initial_vague_answer_count=0,
        vague_answer_count_after_followups=0,
        specificity_improvement_rate=None,
    )
    flags = {
        "h1_voice_more_detailed_support_flag": None,
        "h2_followups_improve_quality_support_flag": None,
        "h3_completion_success_flag": None,
        "h4_extraction_useful_support_flag": None,
        "h5_voice_natural_support_flag": None,
        "h6_device_compatibility_support_flag": None,
    }
    return SubmissionBundle(
        submission=sub,
        answers=[],
        extraction=None,
        review=None,
        cohort=None,
        rollup=rollup,
        hypothesis_flags=flags,
    )


# ─────────────────────────────────────────────────────────────────────────────
# PDF / PPTX summaries — unchanged behaviour, read from bundles
# ─────────────────────────────────────────────────────────────────────────────


def _gather_report_data(bundles: list[SubmissionBundle], cohort_name: str, course_name: str) -> dict:
    survey_version = "1.0"
    if bundles:
        survey_version = bundles[0].submission.survey_version or "1.0"

    total = len(bundles)
    scores: list[int] = []
    all_themes: list[str] = []
    all_barriers: list[str] = []
    all_workflows: list[str] = []
    stories: list[str] = []

    for bundle in bundles:
        for a in bundle.answers:
            if a.question_id == "q1_recommend" and a.answer_raw:
                try:
                    scores.append(int(a.answer_raw))
                except ValueError:
                    pass
        ext = bundle.extraction
        if ext:
            if ext.top_themes:
                all_themes.extend(ext.top_themes)
            if ext.barriers:
                all_barriers.extend(ext.barriers)
            if ext.planned_task_or_workflow:
                all_workflows.append(ext.planned_task_or_workflow)
            if ext.success_story_candidate:
                stories.append(ext.success_story_candidate)

    theme_counts: dict[str, int] = {}
    for t in all_themes:
        theme_counts[t] = theme_counts.get(t, 0) + 1
    top_themes = sorted(theme_counts.items(), key=lambda x: x[1], reverse=True)[:6]

    barrier_counts: dict[str, int] = {}
    for b in all_barriers:
        barrier_counts[b] = barrier_counts.get(b, 0) + 1
    top_barriers = sorted(barrier_counts.items(), key=lambda x: x[1], reverse=True)[:6]

    return {
        "cohort_name": cohort_name,
        "course_name": course_name,
        "survey_version": survey_version,
        "total_responses": total,
        "avg_recommend": round(sum(scores) / len(scores), 1) if scores else None,
        "top_themes": top_themes,
        "top_barriers": top_barriers,
        "workflows": all_workflows[:10],
        "stories": stories[:6],
    }


def generate_summary_pdf(bundles: list[SubmissionBundle], cohort_name: str = "", course_name: str = "") -> bytes:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import inch

    data = _gather_report_data(bundles, cohort_name, course_name)
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("CustomTitle", parent=styles["Title"], textColor=HexColor("#124D8F"), fontSize=20)
    heading_style = ParagraphStyle("CustomHeading", parent=styles["Heading2"], textColor=HexColor("#124D8F"), fontSize=14)

    elements: list = []
    elements.append(Paragraph("InnovateUS Feedback Summary", title_style))
    elements.append(Spacer(1, 12))
    elements.append(Paragraph(f"Cohort: {data['cohort_name'] or 'All'}", styles["Normal"]))
    if data["course_name"]:
        elements.append(Paragraph(f"Course: {data['course_name']} | Survey Version: {data['survey_version']}", styles["Normal"]))
    elements.append(Spacer(1, 20))
    elements.append(Paragraph("Key Metrics", heading_style))
    elements.append(Spacer(1, 8))
    metrics_text = f"Total Responses: {data['total_responses']}"
    if data["avg_recommend"]:
        metrics_text += f" | Avg. Recommendation Score: {data['avg_recommend']}/10"
    elements.append(Paragraph(metrics_text, styles["Normal"]))
    elements.append(Spacer(1, 16))

    if data["top_themes"]:
        elements.append(Paragraph("Top Themes", heading_style))
        elements.append(Spacer(1, 8))
        for theme, count in data["top_themes"]:
            elements.append(Paragraph(f"• {theme} ({count})", styles["Normal"]))
        elements.append(Spacer(1, 16))

    if data["top_barriers"]:
        elements.append(Paragraph("Top Barriers", heading_style))
        elements.append(Spacer(1, 8))
        for barrier, count in data["top_barriers"]:
            elements.append(Paragraph(f"• {barrier} ({count})", styles["Normal"]))
        elements.append(Spacer(1, 16))

    if data["stories"]:
        elements.append(Paragraph("Success Stories", heading_style))
        elements.append(Spacer(1, 8))
        for i, story in enumerate(data["stories"], 1):
            elements.append(Paragraph(f'{i}. "{story}"', styles["Normal"]))
            elements.append(Spacer(1, 4))
        elements.append(Spacer(1, 16))

    if data["workflows"]:
        elements.append(Paragraph("Planned Workflows", heading_style))
        elements.append(Spacer(1, 8))
        for wf in data["workflows"][:8]:
            elements.append(Paragraph(f"• {wf}", styles["Normal"]))

    doc.build(elements)
    return buffer.getvalue()


def generate_summary_pptx(bundles: list[SubmissionBundle], cohort_name: str = "", course_name: str = "") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    data = _gather_report_data(bundles, cohort_name, course_name)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    brand_blue = RGBColor(0x12, 0x4D, 0x8F)

    def add_title_slide(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = brand_blue
        if subtitle and slide.placeholders[1]:
            slide.placeholders[1].text = subtitle

    def add_content_slide(title: str, bullets: list[str]):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = brand_blue
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.paragraphs[0].text = bullet
                tf.paragraphs[0].font.size = Pt(18)
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)

    cohort_label = data["cohort_name"] or "All Cohorts"
    subtitle = cohort_label
    if data["course_name"]:
        subtitle += f"\nCourse: {data['course_name']} | Survey v{data['survey_version']}"
    add_title_slide("InnovateUS Feedback Summary", subtitle)

    metrics_bullets = [f"Total Responses: {data['total_responses']}"]
    if data["avg_recommend"]:
        metrics_bullets.append(f"Average Recommendation Score: {data['avg_recommend']}/10")
    add_content_slide("Participation Metrics", metrics_bullets)

    workflows = data["workflows"][:6] or ["No data yet"]
    add_content_slide("What Learners Plan to Try", workflows)

    barrier_bullets = [f"{b} ({c})" for b, c in data["top_barriers"]] or ["No barriers reported"]
    add_content_slide("Barriers Identified", barrier_bullets)

    story_bullets = [f'"{s}"' for s in data["stories"][:4]] or ["No success stories yet"]
    add_content_slide("Success Stories", story_bullets)

    theme_bullets = [f"{t} ({c})" for t, c in data["top_themes"]] or ["No themes extracted"]
    add_content_slide("Top Themes", theme_bullets)

    add_content_slide("Recommendations", [
        "Continue iterating on course content based on feedback",
        "Address identified barriers in future cohorts",
        "Highlight success stories to promote engagement",
    ])

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
